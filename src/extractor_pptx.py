from __future__ import annotations

import sys
from pathlib import Path

from models import ExtractionResult


def extract_pptx(path: Path, max_chars: int = 8000) -> ExtractionResult:
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        lines: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_lines: list[str] = []

            # Title first
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text_frame.text.strip()
                if title_text:
                    slide_lines.append(f"[슬라이드 {slide_num}] {title_text}")

            # Body text from all shapes
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_lines.append(text)
                # Table content
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if cells:
                            slide_lines.append(" | ".join(cells))

            # Speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_lines.append(f"[노트] {notes_text}")

            lines.extend(slide_lines)

        combined = "\n".join(lines).strip()
        if not combined:
            return ExtractionResult(
                file_type="pptx",
                extraction_status="empty_text",
                notes=["No text extracted from PPTX."],
            )
        return ExtractionResult(
            file_type="pptx",
            extraction_status="success",
            extracted_text=combined,
            text_excerpt=combined[:max_chars],
            page_count=len(prs.slides),
            notes=[f"PPTX text extracted. {len(prs.slides)} slides."],
        )
    except ModuleNotFoundError:
        return ExtractionResult(
            file_type="pptx",
            extraction_status="missing_dependency:python-pptx",
            notes=["python-pptx is not installed. Run: pip install python-pptx"],
        )
    except Exception as exc:
        if sys.platform != "win32":
            return ExtractionResult(
                file_type="pptx",
                extraction_status=f"error:{exc}",
                notes=["PPTX extraction failed."],
            )
        # Windows: MIP 보호 파일일 수 있으므로 PowerPoint COM으로 재시도
        return _extract_pptx_via_com(path, max_chars, fallback_error=str(exc))


def _extract_pptx_via_com(path: Path, max_chars: int, fallback_error: str) -> ExtractionResult:
    """MIP-보호 PPTX를 PowerPoint COM으로 열어 텍스트 추출."""
    try:
        import win32com.client  # type: ignore
    except ImportError:
        return ExtractionResult(
            file_type="pptx",
            extraction_status=f"error:{fallback_error}",
            notes=["PPTX extraction failed. pywin32 not installed for COM fallback."],
        )

    ppt = None
    try:
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = False
        pres = ppt.Presentations.Open(
            str(path.resolve()),
            ReadOnly=True,
            Untitled=False,
            WithWindow=False,
        )
        lines: list[str] = []
        slide_count = pres.Slides.Count
        for i in range(1, slide_count + 1):
            slide = pres.Slides(i)
            for j in range(1, slide.Shapes.Count + 1):
                try:
                    shape = slide.Shapes(j)
                    if shape.HasTextFrame:
                        text = shape.TextFrame.TextRange.Text.strip()
                        if text:
                            lines.append(text)
                except Exception:
                    pass
        pres.Close()
        combined = "\n".join(lines).strip()
        if not combined:
            return ExtractionResult(
                file_type="pptx",
                extraction_status="empty_text",
                notes=["COM fallback: No text extracted from PPTX."],
            )
        return ExtractionResult(
            file_type="pptx",
            extraction_status="success",
            extracted_text=combined,
            text_excerpt=combined[:max_chars],
            page_count=slide_count,
            notes=[f"COM fallback (MIP): PPTX text extracted. {slide_count} slides."],
        )
    except Exception as com_exc:
        return ExtractionResult(
            file_type="pptx",
            extraction_status=f"error:{fallback_error}",
            notes=[f"PPTX extraction failed. COM fallback also failed: {com_exc}"],
        )
    finally:
        if ppt is not None:
            try:
                ppt.Quit()
            except Exception:
                pass
